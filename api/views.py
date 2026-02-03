from django.contrib.auth.models import User
from django.contrib.auth import authenticate, login, logout
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
from .models import Student, Course, Graph, Task
import json
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import io
import datetime
from .ai_service import extract_course_structure, generate_smart_tasks, find_cross_connections, refine_syllabus_with_doubao

@csrf_exempt
def upload_course_view(request):
    if request.method == 'POST':
        try:
            if 'file' not in request.FILES:
                return JsonResponse({'status': 'error', 'message': 'No file uploaded'})
            
            file = request.FILES['file']
            file_name = file.name
            file_content = file.read()
            
            # Simple parsing logic
            text_content = ""
            if file_name.endswith('.pdf'):
                try:
                    reader = PdfReader(io.BytesIO(file_content))
                    for page in reader.pages:
                        text_content += page.extract_text()
                except:
                    text_content = "PDF Parsing Failed"
            elif file_name.endswith('.docx'):
                try:
                    doc = DocxDocument(io.BytesIO(file_content))
                    for para in doc.paragraphs:
                        text_content += para.text + "\n"
                except:
                    text_content = "Docx Parsing Failed"
            elif file_name.endswith('.pptx'):
                try:
                    prs = Presentation(io.BytesIO(file_content))
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text_content += shape.text + "\n"
                except Exception as e:
                    print(f"PPTX Parsing Error: {e}")
                    text_content = "PPTX Parsing Failed"
            else:
                text_content = file_content.decode('utf-8', errors='ignore')

            if not request.user.is_authenticated:
                 return JsonResponse({'status': 'error', 'message': 'Not authenticated'})
            
            student = Student.objects.get(username=request.user.username)
            
            # --- 🚀 NEW: Doubao Refinement Pipeline ---
            print(f"DEBUG: Starting Doubao refinement for {file_name}")
            refined_content = refine_syllabus_with_doubao(text_content)
            
            course = Course(
                name=file_name.split('.')[0],
                outline_text=text_content[:15000], 
                refined_text=refined_content,
                owner=student,
                icon="dumpling"
            )
            course.save()
            
            return JsonResponse({'status': 'success', 'course_id': str(course.id)})
        except Exception as e:
             return JsonResponse({'status': 'error', 'message': str(e)})
    return JsonResponse({'status': 'error', 'message': 'Method not allowed'}, status=405)

@csrf_exempt
def set_thinking_type_view(request):
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            thinking_type = data.get('thinking_type')
            
            if not request.user.is_authenticated:
                 return JsonResponse({'status': 'error', 'message': 'Not authenticated'})
                 
            student = Student.objects.get(username=request.user.username)
            student.thinking_type = thinking_type
            student.save()
            
            return JsonResponse({'status': 'success'})
        except Exception as e:
            return JsonResponse({'status': 'error', 'message': str(e)})
    return JsonResponse({'status': 'error', 'message': 'Method not allowed'}, status=405)

@csrf_exempt
def generate_tasks_view(request):
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            count = int(data.get('count', 3))
            
            if not request.user.is_authenticated:
                 return JsonResponse({'status': 'error', 'message': 'Not authenticated'})
            
            student = Student.objects.get(username=request.user.username)
            
            course = Course.objects.filter(owner=student).order_by('-created_at').first()
            if not course:
                return JsonResponse({'status': 'error', 'message': 'No course found'})
            
            # 1. AI Analysis (Structure Extraction)
            print(f"DEBUG: Starting task generation for: {course.name}")
            print(f"DEBUG: Thinking Type: {student.thinking_type}")
            
            # Use refined text if available, otherwise fallback to raw text
            analysis_source = course.refined_text if course.refined_text else course.outline_text
            print(f"DEBUG: Using {'REFINED' if course.refined_text else 'RAW'} text for analysis")
            
            ai_data = extract_course_structure(analysis_source, student.thinking_type)
            print(f"DEBUG: AI Data Received: {bool(ai_data)}")
            
            nodes = []
            edges = []
            tasks_content = []
            
            if ai_data:
                # Use AI Data
                print("DEBUG: Processing AI Data...")
                raw_nodes = ai_data.get('nodes', [])
                raw_edges = ai_data.get('edges', [])
                concepts = ai_data.get('concepts', [])
                
                print(f"DEBUG: Found {len(raw_nodes)} nodes and {len(raw_edges)} edges")
                
                # Ensure all IDs are strings to prevent Vis.js mismatch
                for n in raw_nodes:
                    n['id'] = str(n.get('id', ''))
                    nodes.append(n)
                
                for e in raw_edges:
                    e['from'] = str(e.get('from', ''))
                    e['to'] = str(e.get('to', ''))
                    edges.append(e)
                
                # Update Course with concepts
                course.extracted_concepts = concepts
                course.save()
                
                # 2. Cross-Course Connections
                # Use mongoengine syntax correctly
                other_courses = Course.objects.filter(owner=student, id__ne=course.id)
                others_data = [{'name': c.name, 'concepts': c.extracted_concepts} for c in other_courses if c.extracted_concepts]
                
                if others_data:
                    print(f"DEBUG: Checking cross-connections with {len(others_data)} other courses")
                    cross_links = find_cross_connections(course.name, concepts, others_data)
                    print(f"DEBUG: Found {len(cross_links)} cross links")
                    
                    # Add cross-links to graph
                    for link in cross_links:
                        # Create a special node for the external concept
                        ext_node_id = f"ext_{link['to_course']}_{link['to_concept']}"
                        nodes.append({
                            'id': ext_node_id,
                            'label': f"{link['to_concept']} ({link['to_course']})",
                            'shape': 'diamond', # Different shape for external
                            'color': '#81D4FA', # Blue for external
                            'level': 2, # Default level for hierarchical layout
                            'title': f"From course: {link['to_course']}\nReason: {link.get('reason', '')}"
                        })
                        
                        # Find local node to connect to (fuzzy match)
                        local_node_id = None
                        for n in nodes:
                            if n['label'].lower() in link['from_concept'].lower() or link['from_concept'].lower() in n['label'].lower():
                                local_node_id = n['id']
                                break
                        
                        if local_node_id:
                            edges.append({
                                'from': local_node_id,
                                'to': ext_node_id,
                                'dashes': True, # Dashed line for cross-link
                                'label': 'Related',
                                'title': link.get('reason', 'Cross-course connection')
                            })

                # 3. AI Task Generation
                print("DEBUG: Generating smart tasks...")
                ai_tasks = generate_smart_tasks(course.name, nodes, count)
                if ai_tasks and 'tasks' in ai_tasks:
                    tasks_content = ai_tasks['tasks']
                    print(f"DEBUG: Generated {len(tasks_content)} tasks")
            else:
                print("DEBUG: AI Data was None, falling back to mock data.")
            
            # Fallback if AI fails or returns empty
            if not tasks_content:
                print("DEBUG: Falling back to mock tasks.")
                tasks_content = [f"Cook {course.name} - Step {i+1}" for i in range(count)]
            
            if not nodes:
                print("DEBUG: Falling back to mock nodes/edges.")
                nodes = [
                    {'id': '1', 'label': course.name, 'shape': 'box', 'color': '#FFD54F', 'level': 0},
                    {'id': '2', 'label': 'Preparation', 'shape': 'dot', 'color': '#FFAB91', 'level': 1},
                    {'id': '3', 'label': 'Core Ingredients', 'shape': 'dot', 'color': '#FFAB91', 'level': 1},
                ]
                edges = [{'from': '1', 'to': '2'}, {'from': '1', 'to': '3'}]

            # Save Tasks
            tasks_data = []
            for i, content in enumerate(tasks_content):
                task = Task(
                    content=content,
                    course=course,
                    owner=student,
                    status="pending"
                )
                task.save()
                tasks_data.append(str(task.id))
                
            # Save Graph
            graph = Graph(
                course=course,
                nodes=nodes,
                edges=edges,
                owner=student
            )
            graph.save()
            
            return JsonResponse({'status': 'success', 'graph_id': str(graph.id), 'task_ids': tasks_data})
        except Exception as e:
            import traceback
            traceback.print_exc()
            return JsonResponse({'status': 'error', 'message': str(e)})
    return JsonResponse({'status': 'error', 'message': 'Method not allowed'}, status=405)


@csrf_exempt
def register_view(request):
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            username = data.get('username')
            password = data.get('password')
            major = data.get('major')
            
            if not username or not password:
                return JsonResponse({'status': 'error', 'message': 'Username and password required'})

            if User.objects.filter(username=username).exists():
                return JsonResponse({'status': 'error', 'message': 'Username exists'})
            
            user = User.objects.create_user(username=username, password=password)
            
            # Create Mongo Student
            student = Student(username=username)
            student.save()
            
            # Piggy complaint logic (Simple stub)
            complaint = f"Wow, {major}? Sounds tasty but tough!"
            
            login(request, user)
            return JsonResponse({'status': 'success', 'complaint': complaint})
        except Exception as e:
            return JsonResponse({'status': 'error', 'message': str(e)})
    return JsonResponse({'status': 'error', 'message': 'Method not allowed'}, status=405)

@csrf_exempt
def login_view(request):
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            username = data.get('username')
            password = data.get('password')
            
            user = authenticate(request, username=username, password=password)
            if user is not None:
                login(request, user)
                try:
                    student = Student.objects.get(username=username)
                    return JsonResponse({
                        'status': 'success', 
                        'username': username,
                        'thinking_type': student.thinking_type
                    })
                except Student.DoesNotExist:
                    # Auto-create if missing (legacy/error recovery)
                    student = Student(username=username)
                    student.save()
                    return JsonResponse({
                        'status': 'success', 
                        'username': username,
                        'thinking_type': student.thinking_type
                    })
            else:
                return JsonResponse({'status': 'error', 'message': 'Invalid credentials'})
        except Exception as e:
            return JsonResponse({'status': 'error', 'message': str(e)})
    return JsonResponse({'status': 'error', 'message': 'Method not allowed'}, status=405)

def check_auth_view(request):
    if request.user.is_authenticated:
        try:
            student = Student.objects.get(username=request.user.username)
            return JsonResponse({
                'is_authenticated': True,
                'username': request.user.username,
                'thinking_type': student.thinking_type
            })
        except Student.DoesNotExist:
            return JsonResponse({'is_authenticated': True, 'username': request.user.username, 'thinking_type': None})
    return JsonResponse({'is_authenticated': False})

def logout_view(request):
    logout(request)
    return JsonResponse({'status': 'success'})

@csrf_exempt
def get_task_details_view(request):
    task_id = request.GET.get('id')
    try:
        task = Task.objects.get(id=task_id)
        return JsonResponse({
            'status': 'success',
            'task': {'id': str(task.id), 'content': task.content, 'status': task.status, 'course_name': task.course.name, 'course_icon': task.course.icon}
        })
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)})

@csrf_exempt
def get_dashboard_data_view(request):
    if not request.user.is_authenticated:
         return JsonResponse({'status': 'error', 'message': 'Not authenticated'})
    
    student = Student.objects.get(username=request.user.username)
    course = Course.objects.filter(owner=student).order_by('-created_at').first()
    if not course:
         return JsonResponse({'status': 'error', 'message': 'No course'})
         
    graph = Graph.objects.filter(course=course).order_by('-created_at').first()
    # Flexible date filter or just take latest batch
    tasks = Task.objects.filter(course=course).order_by('-date')[:5] # Simple fetch latest 5 for demo
    
    tasks_data = [{'id': str(t.id), 'content': t.content, 'status': t.status, 'is_completed': t.is_completed} for t in tasks]
    
    return JsonResponse({
        'status': 'success',
        'thinking_type': student.thinking_type,
        'graph': {
            'nodes': graph.nodes if graph else [],
            'edges': graph.edges if graph else []
        },
        'tasks': tasks_data
    })

@csrf_exempt
def complete_task_view(request):
    if request.method == 'POST':
        data = json.loads(request.body)
        task_id = data.get('task_id')
        status = data.get('status', 'completed')
        
        try:
            task = Task.objects.get(id=task_id)
            task.status = status
            task.is_completed = (status == 'completed')
            task.save()
            
            # Check if all relevant tasks are done
            # For demo, check the tasks returned in dashboard (latest batch)
            all_tasks = Task.objects.filter(course=task.course).order_by('-date')[:5]
            all_done = all([t.is_completed or t.status == 'skipped' for t in all_tasks])
            
            return JsonResponse({'status': 'success', 'all_done': all_done})
        except Exception as e:
            return JsonResponse({'status': 'error', 'message': str(e)})
    return JsonResponse({'status': 'error'})

@csrf_exempt
def get_results_view(request):
    if not request.user.is_authenticated:
         return JsonResponse({'status': 'error', 'message': 'Not authenticated'})
    student = Student.objects.get(username=request.user.username)
    
    # Simple calculation
    tasks = Task.objects.filter(owner=student, is_completed=True).order_by('-date')[:5]
    earned = len(tasks)
    
    return JsonResponse({'status': 'success', 'carrots': earned, 'total_carrots': student.carrots + earned})
